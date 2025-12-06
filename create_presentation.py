#!/usr/bin/env python3
"""
MIDOR-ETHYDCO Integration Presentation Generator
Creates a professional PowerPoint presentation with visuals and text
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import nsmap
from pptx.oxml import parse_xml
import plotly.graph_objects as go
import plotly.io as pio
from io import BytesIO

# Color scheme
COLORS = {
    'primary': RGBColor(14, 165, 233),      # #0ea5e9 - Blue
    'secondary': RGBColor(6, 182, 212),     # #06b6d4 - Cyan
    'accent': RGBColor(245, 158, 11),       # #f59e0b - Orange
    'success': RGBColor(34, 197, 94),       # #22c55e - Green
    'danger': RGBColor(239, 68, 68),        # #ef4444 - Red
    'dark': RGBColor(15, 23, 42),           # #0f172a - Dark blue
    'dark_light': RGBColor(30, 41, 59),     # #1e293b - Dark slate
    'white': RGBColor(255, 255, 255),
    'gray': RGBColor(148, 163, 184),        # #94a3b8
    'light': RGBColor(241, 245, 249),       # #f1f5f9
}

def set_slide_background(slide, color):
    """Set slide background color."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_title_slide(prs):
    """Create title slide."""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, COLORS['dark'])

    # Add gradient overlay shape
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS['dark']
    shape.line.fill.background()

    # Main title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "MIDOR-ETHYDCO Integration"
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(9), Inches(0.8))
    tf = subtitle_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Petrochemical Integration Analysis"
    p.font.size = Pt(28)
    p.font.color.rgb = COLORS['secondary']
    p.alignment = PP_ALIGN.CENTER

    # Value highlight
    value_box = slide.shapes.add_textbox(Inches(0.5), Inches(5), Inches(9), Inches(1))
    tf = value_box.text_frame
    p = tf.paragraphs[0]
    p.text = "$196 Million/Year Net Value"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = COLORS['accent']
    p.alignment = PP_ALIGN.CENTER

    # Footer
    footer_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(9), Inches(0.5))
    tf = footer_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Middle East Oil Refinery (MIDOR) & Egyptian Ethylene and Derivatives Company (ETHYDCO)"
    p.font.size = Pt(14)
    p.font.color.rgb = COLORS['gray']
    p.alignment = PP_ALIGN.CENTER

    return slide

def add_executive_summary(prs):
    """Create executive summary slide."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, COLORS['dark'])

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Executive Summary"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']

    # Key message box
    msg_shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.2), Inches(9), Inches(1.2))
    msg_shape.fill.solid()
    msg_shape.fill.fore_color.rgb = COLORS['dark_light']
    msg_shape.line.color.rgb = COLORS['secondary']
    msg_shape.line.width = Pt(2)

    tf = msg_shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Strategic integration between MIDOR refinery and ETHYDCO petrochemical complex creates significant value through gas recovery, hydrogen utilization, and methanol production pathways."
    p.font.size = Pt(16)
    p.font.color.rgb = COLORS['light']
    p.alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].space_before = Pt(15)

    # Three KPI cards
    kpis = [
        ("$196M", "Total Net Value", "Per Year", COLORS['secondary']),
        ("$100M", "Phase 1+2", "Gas Recovery", COLORS['primary']),
        ("$96M", "Phase 3+4", "Methanol & MTO", COLORS['accent']),
    ]

    for i, (value, title, subtitle, color) in enumerate(kpis):
        x = Inches(0.5 + i * 3.1)

        # Card background
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(2.6), Inches(2.9), Inches(1.8))
        card.fill.solid()
        card.fill.fore_color.rgb = COLORS['dark_light']
        card.line.color.rgb = color
        card.line.width = Pt(3)

        # Value
        val_box = slide.shapes.add_textbox(x, Inches(2.8), Inches(2.9), Inches(0.8))
        tf = val_box.text_frame
        p = tf.paragraphs[0]
        p.text = value
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = color
        p.alignment = PP_ALIGN.CENTER

        # Title
        title_box = slide.shapes.add_textbox(x, Inches(3.5), Inches(2.9), Inches(0.4))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = COLORS['white']
        p.alignment = PP_ALIGN.CENTER

        # Subtitle
        sub_box = slide.shapes.add_textbox(x, Inches(3.85), Inches(2.9), Inches(0.4))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(12)
        p.font.color.rgb = COLORS['gray']
        p.alignment = PP_ALIGN.CENTER

    # Key benefits section
    benefits_title = slide.shapes.add_textbox(Inches(0.5), Inches(4.6), Inches(9), Inches(0.5))
    tf = benefits_title.text_frame
    p = tf.paragraphs[0]
    p.text = "Key Benefits"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = COLORS['secondary']

    benefits = [
        "Eliminates flare gas waste - converts to valuable products",
        "Provides 49-71% of ETHYDCO's ethane feedstock requirements",
        "Creates new revenue streams from hydrogen and methanol",
        "Reduces environmental impact through emission reduction"
    ]

    for i, benefit in enumerate(benefits):
        y = Inches(5.1 + i * 0.4)
        bullet_box = slide.shapes.add_textbox(Inches(0.7), y, Inches(8.5), Inches(0.4))
        tf = bullet_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"✓  {benefit}"
        p.font.size = Pt(14)
        p.font.color.rgb = COLORS['light']

    return slide

def add_problem_statement(prs):
    """Create problem statement slide."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, COLORS['dark'])

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "The Opportunity"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']

    # Two columns layout
    # Left column - MIDOR
    midor_title = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(4.3), Inches(0.5))
    tf = midor_title.text_frame
    p = tf.paragraphs[0]
    p.text = "MIDOR Refinery"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = COLORS['primary']

    midor_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.7), Inches(4.3), Inches(2.5))
    midor_box.fill.solid()
    midor_box.fill.fore_color.rgb = COLORS['dark_light']
    midor_box.line.color.rgb = COLORS['primary']
    midor_box.line.width = Pt(2)

    midor_issues = [
        "• Flaring valuable gases",
        "• 56,000+ t/y of recoverable gases",
        "• Lost hydrogen, LPG, ethane",
        "• Environmental concerns",
        "• Wasted economic potential"
    ]

    for i, issue in enumerate(midor_issues):
        issue_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.9 + i * 0.45), Inches(4), Inches(0.4))
        tf = issue_box.text_frame
        p = tf.paragraphs[0]
        p.text = issue
        p.font.size = Pt(14)
        p.font.color.rgb = COLORS['light']

    # Right column - ETHYDCO
    ethydco_title = slide.shapes.add_textbox(Inches(5.2), Inches(1.2), Inches(4.3), Inches(0.5))
    tf = ethydco_title.text_frame
    p = tf.paragraphs[0]
    p.text = "ETHYDCO Complex"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = COLORS['accent']

    ethydco_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.2), Inches(1.7), Inches(4.3), Inches(2.5))
    ethydco_box.fill.solid()
    ethydco_box.fill.fore_color.rgb = COLORS['dark_light']
    ethydco_box.line.color.rgb = COLORS['accent']
    ethydco_box.line.width = Pt(2)

    ethydco_issues = [
        "• Needs ethane feedstock",
        "• 84-122 kt/y C2 demand",
        "• Dependent on imports",
        "• Supply chain risks",
        "• Cost pressures"
    ]

    for i, issue in enumerate(ethydco_issues):
        issue_box = slide.shapes.add_textbox(Inches(5.4), Inches(1.9 + i * 0.45), Inches(4), Inches(0.4))
        tf = issue_box.text_frame
        p = tf.paragraphs[0]
        p.text = issue
        p.font.size = Pt(14)
        p.font.color.rgb = COLORS['light']

    # Arrow in middle
    arrow_box = slide.shapes.add_textbox(Inches(4.3), Inches(2.7), Inches(1.4), Inches(0.5))
    tf = arrow_box.text_frame
    p = tf.paragraphs[0]
    p.text = "↔"
    p.font.size = Pt(48)
    p.font.color.rgb = COLORS['success']
    p.alignment = PP_ALIGN.CENTER

    # Solution box
    solution_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(4.5), Inches(9), Inches(2))
    solution_box.fill.solid()
    solution_box.fill.fore_color.rgb = COLORS['dark_light']
    solution_box.line.color.rgb = COLORS['success']
    solution_box.line.width = Pt(3)

    sol_title = slide.shapes.add_textbox(Inches(0.7), Inches(4.7), Inches(8.6), Inches(0.5))
    tf = sol_title.text_frame
    p = tf.paragraphs[0]
    p.text = "The Solution: Strategic Integration"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = COLORS['success']

    solution_text = "By integrating MIDOR's off-gas streams with ETHYDCO's feedstock needs, we create a symbiotic relationship that transforms waste into value. MIDOR's flare gases become ETHYDCO's feedstock, while recovered hydrogen enables methanol production for gasoline blending."

    sol_desc = slide.shapes.add_textbox(Inches(0.7), Inches(5.2), Inches(8.6), Inches(1.2))
    tf = sol_desc.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = solution_text
    p.font.size = Pt(14)
    p.font.color.rgb = COLORS['light']

    return slide

def add_integration_phases(prs):
    """Create integration phases overview slide."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, COLORS['dark'])

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Integration Phases"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']

    phases = [
        {
            "num": "1",
            "title": "LPG & C5+ Recovery",
            "value": "$105M",
            "desc": "Recover propane, butane, and naphtha from all gas streams",
            "color": COLORS['secondary']
        },
        {
            "num": "2",
            "title": "Hydrogen Recovery",
            "value": "$79M",
            "desc": "Extract hydrogen for refinery use and methanol synthesis",
            "color": COLORS['primary']
        },
        {
            "num": "3",
            "title": "Methanol Production",
            "value": "$36M",
            "desc": "Convert CO/CO2 + H2 to methanol for gasoline blending",
            "color": COLORS['accent']
        },
        {
            "num": "4",
            "title": "MTO Conversion",
            "value": "$60M",
            "desc": "Methanol-to-Olefins producing ethylene & propylene",
            "color": COLORS['success']
        }
    ]

    for i, phase in enumerate(phases):
        y = Inches(1.2 + i * 1.4)

        # Phase number circle
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.5), y, Inches(0.7), Inches(0.7))
        circle.fill.solid()
        circle.fill.fore_color.rgb = phase['color']
        circle.line.fill.background()

        num_box = slide.shapes.add_textbox(Inches(0.5), y + Inches(0.1), Inches(0.7), Inches(0.5))
        tf = num_box.text_frame
        p = tf.paragraphs[0]
        p.text = phase['num']
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = COLORS['white']
        p.alignment = PP_ALIGN.CENTER

        # Phase content box
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.4), y, Inches(6.5), Inches(1.1))
        box.fill.solid()
        box.fill.fore_color.rgb = COLORS['dark_light']
        box.line.color.rgb = phase['color']
        box.line.width = Pt(2)

        # Title
        title = slide.shapes.add_textbox(Inches(1.6), y + Inches(0.15), Inches(4), Inches(0.4))
        tf = title.text_frame
        p = tf.paragraphs[0]
        p.text = phase['title']
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = COLORS['white']

        # Description
        desc = slide.shapes.add_textbox(Inches(1.6), y + Inches(0.55), Inches(5), Inches(0.5))
        tf = desc.text_frame
        p = tf.paragraphs[0]
        p.text = phase['desc']
        p.font.size = Pt(12)
        p.font.color.rgb = COLORS['gray']

        # Value
        val = slide.shapes.add_textbox(Inches(6.8), y + Inches(0.25), Inches(1.1), Inches(0.6))
        tf = val.text_frame
        p = tf.paragraphs[0]
        p.text = phase['value']
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = phase['color']
        p.alignment = PP_ALIGN.CENTER

        # Connector line
        if i < 3:
            line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.82), y + Inches(0.75), Inches(0.06), Inches(0.65))
            line.fill.solid()
            line.fill.fore_color.rgb = COLORS['gray']
            line.line.fill.background()

    # Total value box
    total_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.2), Inches(2.5), Inches(1.5), Inches(2))
    total_box.fill.solid()
    total_box.fill.fore_color.rgb = COLORS['dark_light']
    total_box.line.color.rgb = COLORS['white']
    total_box.line.width = Pt(2)

    total_title = slide.shapes.add_textbox(Inches(8.2), Inches(2.7), Inches(1.5), Inches(0.4))
    tf = total_title.text_frame
    p = tf.paragraphs[0]
    p.text = "TOTAL"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = COLORS['gray']
    p.alignment = PP_ALIGN.CENTER

    total_val = slide.shapes.add_textbox(Inches(8.2), Inches(3.1), Inches(1.5), Inches(0.6))
    tf = total_val.text_frame
    p = tf.paragraphs[0]
    p.text = "$196M"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    p.alignment = PP_ALIGN.CENTER

    total_sub = slide.shapes.add_textbox(Inches(8.2), Inches(3.6), Inches(1.5), Inches(0.4))
    tf = total_sub.text_frame
    p = tf.paragraphs[0]
    p.text = "Net/Year"
    p.font.size = Pt(11)
    p.font.color.rgb = COLORS['gray']
    p.alignment = PP_ALIGN.CENTER

    return slide

def add_product_values(prs):
    """Create product values slide."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, COLORS['dark'])

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Annual Product Values"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']

    products = [
        ("LPG (C3+C4)", "$91.5M", "125,509 t/y", COLORS['secondary'], 0.915),
        ("Hydrogen (H2)", "$78.9M", "39,445 t/y", COLORS['primary'], 0.789),
        ("Propylene (MTO)", "$44.3M", "60,703 t/y", COLORS['accent'], 0.443),
        ("Methanol Blend", "$35.8M", "79,540 t/y", COLORS['accent'], 0.358),
        ("Ethane (C2)", "$23.6M", "59,005 t/y", COLORS['primary'], 0.236),
        ("Ethylene (MTO)", "$16.2M", "40,468 t/y", COLORS['accent'], 0.162),
        ("Naphtha (C5+)", "$13.5M", "21,733 t/y", COLORS['secondary'], 0.135),
    ]

    max_width = 7.5  # inches for 100%

    for i, (name, value, qty, color, pct) in enumerate(products):
        y = Inches(1.1 + i * 0.75)

        # Product name
        name_box = slide.shapes.add_textbox(Inches(0.5), y, Inches(2.2), Inches(0.4))
        tf = name_box.text_frame
        p = tf.paragraphs[0]
        p.text = name
        p.font.size = Pt(13)
        p.font.color.rgb = COLORS['light']

        # Bar background
        bar_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.7), y + Inches(0.05), Inches(max_width), Inches(0.35))
        bar_bg.fill.solid()
        bar_bg.fill.fore_color.rgb = COLORS['dark_light']
        bar_bg.line.fill.background()

        # Bar fill
        bar_width = max_width * pct
        bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.7), y + Inches(0.05), Inches(bar_width), Inches(0.35))
        bar.fill.solid()
        bar.fill.fore_color.rgb = color
        bar.line.fill.background()

        # Value label
        val_box = slide.shapes.add_textbox(Inches(2.7 + bar_width + 0.1), y, Inches(1), Inches(0.4))
        tf = val_box.text_frame
        p = tf.paragraphs[0]
        p.text = value
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = COLORS['white']

        # Quantity
        qty_box = slide.shapes.add_textbox(Inches(0.5), y + Inches(0.35), Inches(2.2), Inches(0.3))
        tf = qty_box.text_frame
        p = tf.paragraphs[0]
        p.text = qty
        p.font.size = Pt(10)
        p.font.color.rgb = COLORS['gray']

    # Note about costs
    note_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(6.3), Inches(9), Inches(0.6))
    note_box.fill.solid()
    note_box.fill.fore_color.rgb = COLORS['dark_light']
    note_box.line.color.rgb = COLORS['danger']
    note_box.line.width = Pt(2)

    note = slide.shapes.add_textbox(Inches(0.7), Inches(6.4), Inches(8.6), Inches(0.4))
    tf = note.text_frame
    p = tf.paragraphs[0]
    p.text = "Note: Natural Gas makeup cost of $76.2M/year deducted to arrive at net value of $196M"
    p.font.size = Pt(12)
    p.font.color.rgb = COLORS['light']
    p.alignment = PP_ALIGN.CENTER

    return slide

def add_c2_coverage(prs):
    """Create C2 coverage slide."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, COLORS['dark'])

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "ETHYDCO C2 Feed Coverage"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']

    # Subtitle
    sub = slide.shapes.add_textbox(Inches(0.5), Inches(0.9), Inches(9), Inches(0.5))
    tf = sub.text_frame
    p = tf.paragraphs[0]
    p.text = "MIDOR can supply 59,005 t/y of ethane to ETHYDCO"
    p.font.size = Pt(16)
    p.font.color.rgb = COLORS['gray']

    # Two gauge-like displays
    gauges = [
        ("Minimum Demand", "83,600 t/y", "70.6%", COLORS['success']),
        ("Maximum Demand", "121,600 t/y", "48.5%", COLORS['accent'])
    ]

    for i, (title, demand, coverage, color) in enumerate(gauges):
        x = Inches(0.8 + i * 4.7)

        # Card
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.6), Inches(4), Inches(3.2))
        card.fill.solid()
        card.fill.fore_color.rgb = COLORS['dark_light']
        card.line.color.rgb = color
        card.line.width = Pt(2)

        # Title
        t = slide.shapes.add_textbox(x, Inches(1.8), Inches(4), Inches(0.4))
        tf = t.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = COLORS['white']
        p.alignment = PP_ALIGN.CENTER

        # Demand
        d = slide.shapes.add_textbox(x, Inches(2.2), Inches(4), Inches(0.4))
        tf = d.text_frame
        p = tf.paragraphs[0]
        p.text = f"ETHYDCO needs: {demand}"
        p.font.size = Pt(12)
        p.font.color.rgb = COLORS['gray']
        p.alignment = PP_ALIGN.CENTER

        # Big percentage
        pct = slide.shapes.add_textbox(x, Inches(2.8), Inches(4), Inches(1))
        tf = pct.text_frame
        p = tf.paragraphs[0]
        p.text = coverage
        p.font.size = Pt(54)
        p.font.bold = True
        p.font.color.rgb = color
        p.alignment = PP_ALIGN.CENTER

        # Coverage label
        lbl = slide.shapes.add_textbox(x, Inches(3.8), Inches(4), Inches(0.4))
        tf = lbl.text_frame
        p = tf.paragraphs[0]
        p.text = "Coverage"
        p.font.size = Pt(14)
        p.font.color.rgb = COLORS['gray']
        p.alignment = PP_ALIGN.CENTER

    # Key insight
    insight_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(5.1), Inches(9), Inches(1.5))
    insight_box.fill.solid()
    insight_box.fill.fore_color.rgb = COLORS['dark_light']
    insight_box.line.color.rgb = COLORS['secondary']
    insight_box.line.width = Pt(2)

    insight_title = slide.shapes.add_textbox(Inches(0.7), Inches(5.25), Inches(8.6), Inches(0.4))
    tf = insight_title.text_frame
    p = tf.paragraphs[0]
    p.text = "Key Insight"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLORS['secondary']

    insight_text = slide.shapes.add_textbox(Inches(0.7), Inches(5.65), Inches(8.6), Inches(0.9))
    tf = insight_text.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "MIDOR's integration can provide nearly half to over two-thirds of ETHYDCO's ethane requirements, significantly reducing import dependency and creating a reliable local supply chain worth $23.6M annually."
    p.font.size = Pt(13)
    p.font.color.rgb = COLORS['light']

    return slide

def add_hydrogen_balance(prs):
    """Create hydrogen balance slide."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, COLORS['dark'])

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Hydrogen Balance for Methanol"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']

    # Three columns showing H2 balance
    columns = [
        ("H2 Available", "39.4K t/y", "From gas recovery", COLORS['success']),
        ("H2 Required", "65.7K t/y", "For full methanol production", COLORS['primary']),
        ("H2 Deficit", "26.2K t/y", "External supply needed", COLORS['danger']),
    ]

    for i, (title, value, desc, color) in enumerate(columns):
        x = Inches(0.5 + i * 3.1)

        # Card
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.2), Inches(2.9), Inches(2.2))
        card.fill.solid()
        card.fill.fore_color.rgb = COLORS['dark_light']
        card.line.color.rgb = color
        card.line.width = Pt(2)

        # Title
        t = slide.shapes.add_textbox(x, Inches(1.4), Inches(2.9), Inches(0.4))
        tf = t.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(14)
        p.font.color.rgb = COLORS['gray']
        p.alignment = PP_ALIGN.CENTER

        # Value
        v = slide.shapes.add_textbox(x, Inches(1.9), Inches(2.9), Inches(0.8))
        tf = v.text_frame
        p = tf.paragraphs[0]
        p.text = value
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = color
        p.alignment = PP_ALIGN.CENTER

        # Description
        d = slide.shapes.add_textbox(x, Inches(2.7), Inches(2.9), Inches(0.5))
        tf = d.text_frame
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(11)
        p.font.color.rgb = COLORS['gray']
        p.alignment = PP_ALIGN.CENTER

    # Utilization highlight
    util_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3), Inches(3.6), Inches(4), Inches(1))
    util_box.fill.solid()
    util_box.fill.fore_color.rgb = COLORS['dark_light']
    util_box.line.color.rgb = COLORS['secondary']
    util_box.line.width = Pt(3)

    util_title = slide.shapes.add_textbox(Inches(3), Inches(3.75), Inches(4), Inches(0.4))
    tf = util_title.text_frame
    p = tf.paragraphs[0]
    p.text = "H2 Utilization Rate"
    p.font.size = Pt(14)
    p.font.color.rgb = COLORS['gray']
    p.alignment = PP_ALIGN.CENTER

    util_val = slide.shapes.add_textbox(Inches(3), Inches(4.05), Inches(4), Inches(0.5))
    tf = util_val.text_frame
    p = tf.paragraphs[0]
    p.text = "60%"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = COLORS['secondary']
    p.alignment = PP_ALIGN.CENTER

    # Methanol allocation section
    alloc_title = slide.shapes.add_textbox(Inches(0.5), Inches(4.8), Inches(9), Inches(0.5))
    tf = alloc_title.text_frame
    p = tf.paragraphs[0]
    p.text = "Methanol Allocation (224,070 t/y Total)"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']

    allocations = [
        ("Gasoline Blending", "79,540 t/y", "35.5%", "$35.8M", COLORS['secondary']),
        ("MTO Conversion", "144,530 t/y", "64.5%", "$60.4M", COLORS['accent']),
    ]

    for i, (name, qty, pct, value, color) in enumerate(allocations):
        y = Inches(5.3 + i * 0.7)

        # Name
        n = slide.shapes.add_textbox(Inches(0.5), y, Inches(2), Inches(0.4))
        tf = n.text_frame
        p = tf.paragraphs[0]
        p.text = name
        p.font.size = Pt(14)
        p.font.color.rgb = COLORS['light']

        # Bar
        bar_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.5), y + Inches(0.05), Inches(5), Inches(0.35))
        bar_bg.fill.solid()
        bar_bg.fill.fore_color.rgb = COLORS['dark_light']
        bar_bg.line.fill.background()

        bar_width = 5 * float(pct.replace('%', '')) / 100
        bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.5), y + Inches(0.05), Inches(bar_width), Inches(0.35))
        bar.fill.solid()
        bar.fill.fore_color.rgb = color
        bar.line.fill.background()

        # Percentage and value
        info = slide.shapes.add_textbox(Inches(7.6), y, Inches(2), Inches(0.4))
        tf = info.text_frame
        p = tf.paragraphs[0]
        p.text = f"{pct} | {value}"
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = color

    return slide

def add_financial_summary(prs):
    """Create financial summary slide."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, COLORS['dark'])

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Financial Summary"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']

    # Table header
    headers = ["Category", "Gross Value", "NG Cost", "Net Value"]
    col_widths = [2.5, 2, 2, 2]

    y = Inches(1.2)
    x_start = Inches(0.75)

    # Header row
    x = x_start
    for i, header in enumerate(headers):
        header_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(col_widths[i]), Inches(0.5))
        header_box.fill.solid()
        header_box.fill.fore_color.rgb = COLORS['secondary']
        header_box.line.fill.background()

        txt = slide.shapes.add_textbox(x, y + Inches(0.1), Inches(col_widths[i]), Inches(0.4))
        tf = txt.text_frame
        p = tf.paragraphs[0]
        p.text = header
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = COLORS['white']
        p.alignment = PP_ALIGN.CENTER

        x += Inches(col_widths[i])

    # Data rows
    rows = [
        ("Phase 1+2: Gas Recovery", "$176.1M", "-$76.2M", "$99.8M"),
        ("Phase 3+4: Methanol & MTO", "$96.2M", "$0", "$96.2M"),
        ("TOTAL", "$272.3M", "-$76.2M", "$196.1M"),
    ]

    for row_idx, row_data in enumerate(rows):
        y += Inches(0.6)
        x = x_start

        is_total = row_idx == len(rows) - 1

        for col_idx, cell in enumerate(row_data):
            cell_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(col_widths[col_idx]), Inches(0.55))
            cell_box.fill.solid()
            cell_box.fill.fore_color.rgb = COLORS['dark_light'] if not is_total else RGBColor(20, 50, 40)
            cell_box.line.color.rgb = COLORS['gray']
            cell_box.line.width = Pt(0.5)

            txt = slide.shapes.add_textbox(x, y + Inches(0.12), Inches(col_widths[col_idx]), Inches(0.4))
            tf = txt.text_frame
            p = tf.paragraphs[0]
            p.text = cell
            p.font.size = Pt(13) if not is_total else Pt(14)
            p.font.bold = is_total or col_idx == 0

            # Color coding
            if col_idx == 2 and "-" in cell:
                p.font.color.rgb = COLORS['danger']
            elif col_idx == 3:
                p.font.color.rgb = COLORS['success']
            else:
                p.font.color.rgb = COLORS['white']

            if col_idx > 0:
                p.alignment = PP_ALIGN.CENTER

            x += Inches(col_widths[col_idx])

    # Value breakdown visual
    breakdown_title = slide.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(9), Inches(0.5))
    tf = breakdown_title.text_frame
    p = tf.paragraphs[0]
    p.text = "Value Composition"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']

    # Stacked bar visualization
    total_width = 8.5
    phase12_width = total_width * (99.8 / 196.1)
    phase34_width = total_width * (96.2 / 196.1)

    # Phase 1+2 bar
    bar1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.75), Inches(4.4), Inches(phase12_width), Inches(0.8))
    bar1.fill.solid()
    bar1.fill.fore_color.rgb = COLORS['primary']
    bar1.line.fill.background()

    bar1_txt = slide.shapes.add_textbox(Inches(0.75), Inches(4.55), Inches(phase12_width), Inches(0.5))
    tf = bar1_txt.text_frame
    p = tf.paragraphs[0]
    p.text = "Phase 1+2: $99.8M (51%)"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    p.alignment = PP_ALIGN.CENTER

    # Phase 3+4 bar
    bar2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.75 + phase12_width), Inches(4.4), Inches(phase34_width), Inches(0.8))
    bar2.fill.solid()
    bar2.fill.fore_color.rgb = COLORS['accent']
    bar2.line.fill.background()

    bar2_txt = slide.shapes.add_textbox(Inches(0.75 + phase12_width), Inches(4.55), Inches(phase34_width), Inches(0.5))
    tf = bar2_txt.text_frame
    p = tf.paragraphs[0]
    p.text = "Phase 3+4: $96.2M (49%)"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    p.alignment = PP_ALIGN.CENTER

    # ROI note
    roi_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(5.5), Inches(9), Inches(1.3))
    roi_box.fill.solid()
    roi_box.fill.fore_color.rgb = COLORS['dark_light']
    roi_box.line.color.rgb = COLORS['success']
    roi_box.line.width = Pt(2)

    roi_title = slide.shapes.add_textbox(Inches(0.7), Inches(5.65), Inches(8.6), Inches(0.4))
    tf = roi_title.text_frame
    p = tf.paragraphs[0]
    p.text = "Investment Considerations"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLORS['success']

    roi_text = slide.shapes.add_textbox(Inches(0.7), Inches(6.0), Inches(8.6), Inches(0.7))
    tf = roi_text.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "• Net annual value of $196M provides strong basis for capital investment\n• Phased implementation reduces initial capital requirements\n• Phase 1+2 can be implemented independently with positive returns"
    p.font.size = Pt(12)
    p.font.color.rgb = COLORS['light']

    return slide

def add_next_steps(prs):
    """Create next steps / recommendations slide."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, COLORS['dark'])

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Recommendations & Next Steps"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']

    steps = [
        {
            "num": "1",
            "title": "Feasibility Study",
            "desc": "Conduct detailed engineering and economic feasibility study for gas recovery infrastructure",
            "color": COLORS['secondary']
        },
        {
            "num": "2",
            "title": "Partnership Agreement",
            "desc": "Establish formal partnership framework between MIDOR and ETHYDCO for feedstock supply",
            "color": COLORS['primary']
        },
        {
            "num": "3",
            "title": "Phase 1 Implementation",
            "desc": "Begin with LPG and hydrogen recovery as quick wins with proven technology",
            "color": COLORS['accent']
        },
        {
            "num": "4",
            "title": "Methanol Unit Planning",
            "desc": "Plan methanol synthesis and MTO units based on Phase 1 performance",
            "color": COLORS['success']
        },
    ]

    for i, step in enumerate(steps):
        y = Inches(1.1 + i * 1.35)

        # Number box
        num_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), y, Inches(0.6), Inches(0.6))
        num_box.fill.solid()
        num_box.fill.fore_color.rgb = step['color']
        num_box.line.fill.background()

        num_txt = slide.shapes.add_textbox(Inches(0.5), y + Inches(0.1), Inches(0.6), Inches(0.4))
        tf = num_txt.text_frame
        p = tf.paragraphs[0]
        p.text = step['num']
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = COLORS['white']
        p.alignment = PP_ALIGN.CENTER

        # Content
        title_txt = slide.shapes.add_textbox(Inches(1.3), y, Inches(8), Inches(0.5))
        tf = title_txt.text_frame
        p = tf.paragraphs[0]
        p.text = step['title']
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = step['color']

        desc_txt = slide.shapes.add_textbox(Inches(1.3), y + Inches(0.45), Inches(8), Inches(0.5))
        tf = desc_txt.text_frame
        p = tf.paragraphs[0]
        p.text = step['desc']
        p.font.size = Pt(13)
        p.font.color.rgb = COLORS['light']

    # Timeline hint
    timeline = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(5.8), Inches(9), Inches(1))
    timeline.fill.solid()
    timeline.fill.fore_color.rgb = COLORS['dark_light']
    timeline.line.color.rgb = COLORS['secondary']
    timeline.line.width = Pt(2)

    tl_title = slide.shapes.add_textbox(Inches(0.7), Inches(5.95), Inches(8.6), Inches(0.4))
    tf = tl_title.text_frame
    p = tf.paragraphs[0]
    p.text = "Suggested Timeline"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = COLORS['secondary']

    tl_text = slide.shapes.add_textbox(Inches(0.7), Inches(6.3), Inches(8.6), Inches(0.4))
    tf = tl_text.text_frame
    p = tf.paragraphs[0]
    p.text = "Phase 1+2: 18-24 months  |  Phase 3+4: 24-36 months after Phase 1+2 completion"
    p.font.size = Pt(12)
    p.font.color.rgb = COLORS['light']

    return slide

def add_conclusion(prs):
    """Create conclusion slide."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, COLORS['dark'])

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Conclusion"
    p.font.size = Pt(42)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']

    # Main message
    msg_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.6), Inches(9), Inches(1.5))
    msg_box.fill.solid()
    msg_box.fill.fore_color.rgb = COLORS['dark_light']
    msg_box.line.color.rgb = COLORS['secondary']
    msg_box.line.width = Pt(3)

    msg = slide.shapes.add_textbox(Inches(0.7), Inches(1.85), Inches(8.6), Inches(1.2))
    tf = msg.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "The MIDOR-ETHYDCO integration represents a transformative opportunity to create $196 million in annual value while strengthening Egypt's petrochemical industry and reducing environmental impact."
    p.font.size = Pt(18)
    p.font.color.rgb = COLORS['light']
    p.alignment = PP_ALIGN.CENTER

    # Key stats row
    stats = [
        ("$196M", "Annual Net Value"),
        ("6", "Gas Streams Utilized"),
        ("7", "Valuable Products"),
        ("60%+", "Feedstock Coverage"),
    ]

    for i, (value, label) in enumerate(stats):
        x = Inches(0.5 + i * 2.4)

        stat_val = slide.shapes.add_textbox(x, Inches(3.5), Inches(2.2), Inches(0.8))
        tf = stat_val.text_frame
        p = tf.paragraphs[0]
        p.text = value
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = COLORS['secondary']
        p.alignment = PP_ALIGN.CENTER

        stat_lbl = slide.shapes.add_textbox(x, Inches(4.2), Inches(2.2), Inches(0.5))
        tf = stat_lbl.text_frame
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(12)
        p.font.color.rgb = COLORS['gray']
        p.alignment = PP_ALIGN.CENTER

    # Call to action
    cta = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.5), Inches(5), Inches(5), Inches(0.8))
    cta.fill.solid()
    cta.fill.fore_color.rgb = COLORS['success']
    cta.line.fill.background()

    cta_txt = slide.shapes.add_textbox(Inches(2.5), Inches(5.15), Inches(5), Inches(0.5))
    tf = cta_txt.text_frame
    p = tf.paragraphs[0]
    p.text = "Ready to Transform Waste into Value"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    p.alignment = PP_ALIGN.CENTER

    # Contact/footer
    footer = slide.shapes.add_textbox(Inches(0.5), Inches(6.3), Inches(9), Inches(0.5))
    tf = footer.text_frame
    p = tf.paragraphs[0]
    p.text = "MIDOR-ETHYDCO Integration Analysis | December 2025"
    p.font.size = Pt(12)
    p.font.color.rgb = COLORS['gray']
    p.alignment = PP_ALIGN.CENTER

    return slide

def main():
    """Generate the presentation."""
    print("Creating MIDOR-ETHYDCO Integration Presentation...")

    # Create presentation with 16:9 aspect ratio
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Add slides
    add_title_slide(prs)
    add_executive_summary(prs)
    add_problem_statement(prs)
    add_integration_phases(prs)
    add_product_values(prs)
    add_c2_coverage(prs)
    add_hydrogen_balance(prs)
    add_financial_summary(prs)
    add_next_steps(prs)
    add_conclusion(prs)

    # Save presentation
    output_path = 'MIDOR_ETHYDCO_Integration_Presentation.pptx'
    prs.save(output_path)

    print(f"\n{'='*60}")
    print("Presentation created successfully!")
    print(f"{'='*60}")
    print(f"\nOutput: {output_path}")
    print(f"\nSlides: 10")
    print("\nContents:")
    print("  1. Title Slide")
    print("  2. Executive Summary")
    print("  3. The Opportunity")
    print("  4. Integration Phases")
    print("  5. Annual Product Values")
    print("  6. ETHYDCO C2 Feed Coverage")
    print("  7. Hydrogen Balance for Methanol")
    print("  8. Financial Summary")
    print("  9. Recommendations & Next Steps")
    print(" 10. Conclusion")

if __name__ == '__main__':
    main()
